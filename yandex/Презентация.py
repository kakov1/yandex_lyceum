from pptx import Presentation
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Некоторые методы модуля random"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'randint'
tf = slide.placeholders[1].text_frame.paragraphs[0]
run = tf.add_run()
run.text = 'Help on method randint in module random:' \
           ' randint(a, b) method of random.Random instanceReturn' \
           ' random integer in range [a, b], including both end points.None'
font = run.font
font.name = 'Courier New'
font.size = Pt(40)

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'choice'
tf = slide.placeholders[1].text_frame.paragraphs[0]
run = tf.add_run()
run.text = 'Help on method choice in module random: choice(seq) method of random.' \
           ' Random instance Choose a random element from a non-empty sequence. ' \
           'None'
font = run.font
font.name = 'Courier New'
font.size = Pt(40)

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'random'
tf = slide.placeholders[1].text_frame.paragraphs[0]
run = tf.add_run()
run.text = 'Help on built-in function random: random()' \
           ' method of random.Random instance random()' \
           ' -> x in the interval [0, 1). ' \
           'None'
font = run.font
font.name = 'Courier New'
font.size = Pt(40)

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'randrange'
tf = slide.placeholders[1].text_frame.paragraphs[0]
run = tf.add_run()
run.text = 'Help on method randrange in module random:' \
           'randrange(start, stop=None, step=1, _int=<class int>) method of random.Random instance' \
           'Choose a random item from range(start, stop[, step]).' \
           'This fixes the problem with randint() which includes the' \
           'endpoint; in Python this is usually not what you want. ' \
           'None'
font = run.font
font.name = 'Courier New'
font.size = Pt(25)

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'uniform'
tf = slide.placeholders[1].text_frame.paragraphs[0]
run = tf.add_run()
run.text = 'Help on method uniform in module random: ' \
           'uniform(a, b) method of random.Random instance' \
           'Get a random number in the range [a, b) or [a, b] depending on rounding. ' \
           'None'
font = run.font
font.name = 'Courier New'
font.size = Pt(40)

prs.save('res.pptx')
